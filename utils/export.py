import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from fpdf import FPDF

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Colours ────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
GOLD   = RGBColor(0xC9, 0xA8, 0x52)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xAE, 0xAE, 0xB2)
DGREY  = RGBColor(0x1D, 0x1D, 0x1F)
YELLOW = RGBColor(0xFF, 0xC2, 0x2B)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
GREEN  = RGBColor(0x2E, 0xC4, 0x6E)
LGREEN = RGBColor(0xF0, 0xFF, 0xF4)
SLATE  = RGBColor(0x3A, 0x3A, 0x3C)


# ── Low-level helpers ──────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill: RGBColor, line=False):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = fill
    else:
        shape.line.fill.background()
    return shape


def _tb(slide, x, y, w, h, text, pt, bold=False,
        color=DGREY, align=PP_ALIGN.LEFT, italic=False):
    """Add a textbox. Text is a plain string; each \\n becomes a new paragraph."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name  = "Arial"
        run.font.size  = Pt(pt)
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def _bullets_tf(slide, x, y, w, h, bullets, pt, color, prefix, spacing_pt=14):
    """Each bullet is a separate paragraph with proper before-spacing."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True

    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing_pt if i > 0 else 0)
        p.space_after  = Pt(4)
        run = p.add_run()
        run.text = f"{prefix}{text}"
        run.font.name  = "Arial"
        run.font.size  = Pt(pt)
        run.font.color.rgb = color
    return tb


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _add_bg_image(slide, img_bytes: bytes) -> bool:
    """Stretch photo to fill the entire slide. Returns True on success."""
    try:
        slide.shapes.add_picture(io.BytesIO(img_bytes), 0, 0, SLIDE_W, SLIDE_H)
        return True
    except Exception:
        return False


def _add_overlay(slide, color: RGBColor, opacity_pct: int):
    """Full-slide semi-transparent colour rectangle. opacity_pct: 100=opaque, 0=invisible."""
    shape = _rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)
    solid = shape._element.spPr.find('.//' + qn('a:solidFill'))
    if solid is not None:
        a = etree.SubElement(solid[0], qn('a:alpha'))
        a.set('val', str(opacity_pct * 1000))


# ══════════════════════════════════════════════════════════════════════════════
#  FORMAL  —  "Executive Split Panel"
#  Left navy sidebar (2.8 in) | Right white content area
# ══════════════════════════════════════════════════════════════════════════════
SIDEBAR_W = Inches(2.9)
MAIN_X    = SIDEBAR_W
MAIN_W    = SLIDE_W - SIDEBAR_W
TITLE_H   = Inches(1.7)   # fixed height for title band — never overlaps bullets


def _formal_cover(prs, sd, img_bytes=None):
    slide = _blank(prs)
    bg_ok = img_bytes and _add_bg_image(slide, img_bytes)
    if bg_ok:
        _add_overlay(slide, NAVY, 65)      # dark navy tint — photo shows subtly
    else:
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)

    # Gold horizontal rule at 48%
    rule_y = Inches(3.6)
    _rect(slide, Inches(0.9), rule_y, SLIDE_W - Inches(1.8), Inches(0.04), GOLD)

    # Title above rule
    _tb(slide, Inches(0.9), Inches(1.2),
        SLIDE_W - Inches(1.8), Inches(2.2),
        sd.get("title", ""), 44, bold=True, color=WHITE)

    # Subtitle below rule
    _tb(slide, Inches(0.9), rule_y + Inches(0.18),
        SLIDE_W - Inches(1.8), Inches(1.2),
        sd.get("subtitle", ""), 22, color=LGREY)

    # Meta bottom-right in gold
    meta = f"{sd.get('subject','')}   ·   {sd.get('grade','')}"
    _tb(slide, Inches(0.9), Inches(6.8),
        SLIDE_W - Inches(1.2), Inches(0.5),
        meta, 13, color=GOLD, align=PP_ALIGN.RIGHT)

    # Decorative gold corner block (bottom-left)
    _rect(slide, 0, Inches(6.5), Inches(0.5), Inches(1.0), GOLD)

    _notes(slide, sd.get("speaker_notes", ""))


def _formal_content(prs, sd, num):
    slide = _blank(prs)

    # ── Sidebar (navy) ─────────────────────────────────────────────────────────
    _rect(slide, 0, 0, SIDEBAR_W, SLIDE_H, NAVY)

    # Large slide number in sidebar
    _tb(slide, Inches(0.15), Inches(0.3), SIDEBAR_W - Inches(0.2), Inches(1.6),
        f"{num:02d}", 88, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Gold separator line in sidebar
    _rect(slide, Inches(0.35), Inches(2.0), SIDEBAR_W - Inches(0.7), Inches(0.04), GOLD)

    # Subject label in sidebar
    subject = sd.get("subject", "")
    _tb(slide, Inches(0.15), Inches(2.2), SIDEBAR_W - Inches(0.2), Inches(0.6),
        subject.upper(), 10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Decorative bottom circle in sidebar
    _rect(slide, Inches(0.8), Inches(6.0), Inches(1.3), Inches(1.3), GOLD)
    _tb(slide, Inches(0.8), Inches(6.1), Inches(1.3), Inches(1.0),
        "◆", 36, color=NAVY, align=PP_ALIGN.CENTER)

    # ── Main area (white) ──────────────────────────────────────────────────────
    _rect(slide, MAIN_X, 0, MAIN_W, SLIDE_H, WHITE)

    # Navy title band (fixed height = TITLE_H — title can never overflow into bullets)
    _rect(slide, MAIN_X, 0, MAIN_W, TITLE_H, NAVY)

    # Title text inside band (padded, word-wrapped, white)
    _tb(slide, MAIN_X + Inches(0.35), Inches(0.18),
        MAIN_W - Inches(0.5), TITLE_H - Inches(0.18),
        sd.get("title", ""), 26, bold=True, color=WHITE)

    # Bullets in white area — starts AFTER title band, never overlaps
    _bullets_tf(slide,
                MAIN_X + Inches(0.35), TITLE_H + Inches(0.25),
                MAIN_W - Inches(0.55), SLIDE_H - TITLE_H - Inches(0.55),
                sd.get("bullets", []), 19, DGREY, "▸  ", spacing_pt=16)

    _notes(slide, sd.get("speaker_notes", ""))


def _formal_summary(prs, sd, num):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)

    # Gold top accent band
    _rect(slide, 0, 0, SLIDE_W, Inches(0.12), GOLD)

    # "SUMMARY" label
    _tb(slide, Inches(0.6), Inches(0.25),
        SLIDE_W - Inches(1.2), Inches(0.5),
        "SUMMARY", 11, bold=True, color=GOLD, italic=False)

    # Title
    _tb(slide, Inches(0.6), Inches(0.85),
        SLIDE_W - Inches(1.2), Inches(1.4),
        sd.get("title", ""), 34, bold=True, color=GOLD)

    # Gold rule
    _rect(slide, Inches(0.6), Inches(2.35), SLIDE_W - Inches(1.2), Inches(0.04), GOLD)

    # Bullets — two columns for visual impact
    bullets = sd.get("bullets", [])
    half = len(bullets) // 2 + len(bullets) % 2
    col_w = (SLIDE_W - Inches(1.4)) / 2

    _bullets_tf(slide, Inches(0.6), Inches(2.55),
                col_w - Inches(0.2), Inches(4.6),
                bullets[:half], 20, WHITE, "▸  ", spacing_pt=20)

    _bullets_tf(slide, Inches(0.6) + col_w + Inches(0.2), Inches(2.55),
                col_w - Inches(0.2), Inches(4.6),
                bullets[half:], 20, WHITE, "▸  ", spacing_pt=20)

    # Slide number
    _tb(slide, Inches(0.4), Inches(7.1), SLIDE_W - Inches(0.8), Inches(0.3),
        str(num), 11, color=GOLD, align=PP_ALIGN.RIGHT)

    _notes(slide, sd.get("speaker_notes", ""))


# ══════════════════════════════════════════════════════════════════════════════
#  FUN  —  "Pop Block"
#  Bold yellow header block (full width) | Dark charcoal body
# ══════════════════════════════════════════════════════════════════════════════
HEADER_H = Inches(2.2)


def _fun_cover(prs, sd, img_bytes=None):
    slide = _blank(prs)
    bg_ok = img_bytes and _add_bg_image(slide, img_bytes)
    if bg_ok:
        _add_overlay(slide, DARK, 55)      # semi-transparent dark — photo glows through
    else:
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK)

    # Yellow hero block (solid — sits on top of photo or dark bg)
    _rect(slide, 0, 0, SLIDE_W, Inches(3.0), YELLOW)

    # Subject pill inside yellow block
    subj = sd.get("subject", "")
    _tb(slide, Inches(0.7), Inches(0.3), Inches(5), Inches(0.6),
        subj.upper(), 12, bold=True, color=DARK)

    # Big title in yellow block — dark text
    _tb(slide, Inches(0.7), Inches(0.8),
        SLIDE_W - Inches(1.4), Inches(1.8),
        sd.get("title", ""), 42, bold=True, color=DARK)

    # Grade badge top-right
    grade = sd.get("grade", "")
    _tb(slide, SLIDE_W - Inches(3.5), Inches(0.25), Inches(3.2), Inches(0.5),
        grade, 13, bold=True, color=DARK, align=PP_ALIGN.RIGHT)

    # Subtitle in dark area (yellow text)
    _tb(slide, Inches(0.7), Inches(3.3),
        SLIDE_W - Inches(1.4), Inches(1.2),
        sd.get("subtitle", ""), 22, color=YELLOW)

    # Decorative star cluster
    _tb(slide, Inches(0.5), Inches(6.5), Inches(5), Inches(0.7),
        "★  ★  ★", 20, color=YELLOW)

    _notes(slide, sd.get("speaker_notes", ""))


def _fun_content(prs, sd, num):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK)

    # Yellow full-width header block
    _rect(slide, 0, 0, SLIDE_W, HEADER_H, YELLOW)

    # Slide number pill (top-right, in yellow block)
    _rect(slide, SLIDE_W - Inches(1.1), Inches(0.15), Inches(0.8), Inches(0.5), DARK)
    _tb(slide, SLIDE_W - Inches(1.1), Inches(0.15), Inches(0.8), Inches(0.5),
        str(num), 16, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

    # Title in yellow block (dark text, large)
    _tb(slide, Inches(0.5), Inches(0.2),
        SLIDE_W - Inches(1.8), HEADER_H - Inches(0.3),
        sd.get("title", ""), 28, bold=True, color=DARK)

    # Bullets in dark area — starts AFTER header block, no overlap possible
    _bullets_tf(slide,
                Inches(0.6), HEADER_H + Inches(0.3),
                SLIDE_W - Inches(1.2), SLIDE_H - HEADER_H - Inches(0.55),
                sd.get("bullets", []), 21, WHITE, "★  ", spacing_pt=18)

    _notes(slide, sd.get("speaker_notes", ""))


def _fun_summary(prs, sd, num):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, YELLOW)

    # Dark top strip
    _rect(slide, 0, 0, SLIDE_W, Inches(0.1), DARK)

    _tb(slide, Inches(0.6), Inches(0.25),
        SLIDE_W - Inches(1.2), Inches(0.5),
        "THAT'S A WRAP!", 13, bold=True, color=DARK, italic=True)

    _tb(slide, Inches(0.6), Inches(0.8),
        SLIDE_W - Inches(1.2), Inches(1.3),
        sd.get("title", ""), 32, bold=True, color=DARK)

    _rect(slide, Inches(0.6), Inches(2.2), SLIDE_W - Inches(1.2), Inches(0.05), DARK)

    _bullets_tf(slide, Inches(0.6), Inches(2.45),
                SLIDE_W - Inches(1.2), Inches(4.7),
                sd.get("bullets", []), 22, DARK, "★  ", spacing_pt=20)

    _tb(slide, Inches(0.4), Inches(7.1), SLIDE_W - Inches(0.8), Inches(0.3),
        str(num), 11, color=DARK, align=PP_ALIGN.RIGHT)

    _notes(slide, sd.get("speaker_notes", ""))


# ══════════════════════════════════════════════════════════════════════════════
#  SIMPLE  —  "Billboard"
#  Green full-width title band | Clean white body, very large text
# ══════════════════════════════════════════════════════════════════════════════
BILLBOARD_H = Inches(1.9)


def _simple_cover(prs, sd, img_bytes=None):
    slide = _blank(prs)
    bg_ok = img_bytes and _add_bg_image(slide, img_bytes)
    if bg_ok:
        _add_overlay(slide, WHITE, 72)     # light white wash — airy, clean look
    else:
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)

    # Green left half panel (solid — covers photo on left, photo shows on right)
    _rect(slide, 0, 0, Inches(4.5), SLIDE_H, GREEN)

    # Big initial letter in left panel
    subj = sd.get("subject", "Education")
    _tb(slide, Inches(0.2), Inches(1.0), Inches(4.1), Inches(4.0),
        subj[0].upper(), 200, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subject label below initial
    _tb(slide, Inches(0.2), Inches(5.5), Inches(4.1), Inches(0.6),
        subj.upper(), 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Title on right side
    _tb(slide, Inches(5.0), Inches(1.5),
        SLIDE_W - Inches(5.5), Inches(2.8),
        sd.get("title", ""), 40, bold=True, color=DGREY)

    # Subtitle on right
    _tb(slide, Inches(5.0), Inches(4.5),
        SLIDE_W - Inches(5.5), Inches(1.4),
        sd.get("subtitle", ""), 22, color=SLATE)

    # Grade bottom-right
    _tb(slide, Inches(5.0), Inches(6.7),
        SLIDE_W - Inches(5.4), Inches(0.5),
        sd.get("grade", ""), 14, color=GREEN, align=PP_ALIGN.RIGHT)

    _notes(slide, sd.get("speaker_notes", ""))


def _simple_content(prs, sd, num):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)

    # Green full-width title band
    _rect(slide, 0, 0, SLIDE_W, BILLBOARD_H, GREEN)

    # Slide number in top-right of green band
    _tb(slide, SLIDE_W - Inches(1.0), Inches(0.1), Inches(0.8), Inches(0.6),
        str(num), 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Title text inside green band — white, padded
    _tb(slide, Inches(0.5), Inches(0.2),
        SLIDE_W - Inches(1.6), BILLBOARD_H - Inches(0.3),
        sd.get("title", ""), 30, bold=True, color=WHITE)

    # Thin green left bar below the band (visual continuity)
    _rect(slide, 0, BILLBOARD_H, Inches(0.18), SLIDE_H - BILLBOARD_H, GREEN)

    # Bullets in white area — starts AFTER green band, NEVER overlaps
    _bullets_tf(slide,
                Inches(0.45), BILLBOARD_H + Inches(0.3),
                SLIDE_W - Inches(0.65), SLIDE_H - BILLBOARD_H - Inches(0.5),
                sd.get("bullets", []), 23, DGREY, "●  ", spacing_pt=20)

    _notes(slide, sd.get("speaker_notes", ""))


def _simple_summary(prs, sd, num):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREEN)

    # Green top band
    _rect(slide, 0, 0, SLIDE_W, Inches(0.2), GREEN)

    _tb(slide, Inches(0.6), Inches(0.4),
        SLIDE_W - Inches(1.2), Inches(0.5),
        "YOU'VE GOT THIS!", 13, bold=True, color=GREEN)

    _tb(slide, Inches(0.6), Inches(0.95),
        SLIDE_W - Inches(1.2), Inches(1.3),
        sd.get("title", ""), 32, bold=True, color=DGREY)

    _rect(slide, Inches(0.6), Inches(2.35), SLIDE_W - Inches(1.2), Inches(0.07), GREEN)

    _bullets_tf(slide, Inches(0.6), Inches(2.6),
                SLIDE_W - Inches(1.2), Inches(4.5),
                sd.get("bullets", []), 24, DGREY, "●  ", spacing_pt=22)

    _tb(slide, Inches(0.4), Inches(7.1), SLIDE_W - Inches(0.8), Inches(0.3),
        str(num), 11, color=GREEN, align=PP_ALIGN.RIGHT)

    _notes(slide, sd.get("speaker_notes", ""))


# ══════════════════════════════════════════════════════════════════════════════
#  Public entry point
# ══════════════════════════════════════════════════════════════════════════════
_LAYOUTS = {
    "Formal": (_formal_cover, _formal_content, _formal_summary),
    "Fun":    (_fun_cover,    _fun_content,    _fun_summary),
    "Simple": (_simple_cover, _simple_content, _simple_summary),
}


def create_presentation(data: dict, img_bytes: bytes = None) -> bytes:
    tone = data.get("tone", "Formal")
    cover_fn, content_fn, summary_fn = _LAYOUTS.get(tone, _LAYOUTS["Formal"])

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    subject = data.get("subject", "")
    grade   = data.get("grade", "")

    for sd in data.get("slides", []):
        sd.setdefault("subject", subject)
        sd.setdefault("grade", grade)
        stype = sd.get("slide_type", "content")
        snum  = sd.get("slide_number", 0)

        if stype == "cover":
            cover_fn(prs, sd, img_bytes=img_bytes)
        elif stype == "summary":
            summary_fn(prs, sd, snum)
        else:
            content_fn(prs, sd, snum)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ══════════════════════════════════════════════════════════════════════════════
# LESSON PLAN PDF  —  Rich design
# ══════════════════════════════════════════════════════════════════════════════

def _safe(text: str) -> str:
    """Replace Unicode chars that break Helvetica (latin-1 only)."""
    _REPLACEMENTS = {
        # Dashes & quotes
        "—": "-", "–": "-",
        "'": "'", "'": "'",
        """: '"', """: '"',
        # Common math/science
        "•": "-",   # bullet
        "²": "2",   # superscript 2
        "³": "3",   # superscript 3
        "…": "...", # ellipsis
        "×": "x",   # multiplication ×
        "÷": "/",   # division ÷
        "±": "+/-", # ±
        "√": "sqrt",# √
        "∑": "sum", # ∑
        "∫": "int", # ∫
        "∞": "inf", # ∞
        # Arrows
        "→": "->",  # →
        "←": "<-",  # ←
        "↑": "^",   # ↑
        "↓": "v",   # ↓
        "⇒": "=>",  # ⇒
        # Degree & temperature
        "°": " deg",# °
        # Greek letters (common in science/math)
        "π": "pi",  # π
        "α": "alpha",# α
        "β": "beta", # β
        "γ": "gamma",# γ
        "δ": "delta",# δ
        "θ": "theta",# θ
        "λ": "lambda",# λ
        "μ": "mu",  # μ
        "σ": "sigma",# σ
        "Ω": "Omega",# Ω
        # Comparisons
        "≤": "<=",  # ≤
        "≥": ">=",  # ≥
        "≠": "!=",  # ≠
        "≈": "~=",  # ≈
        # India-specific
        "₹": "Rs.", # ₹
        # Fractions
        "¼": "1/4", # ¼
        "½": "1/2", # ½
        "¾": "3/4", # ¾
    }
    s = str(text)
    for old, new in _REPLACEMENTS.items():
        s = s.replace(old, new)
    return s.encode("latin-1", errors="replace").decode("latin-1")


PDF_MARGIN = 14
PDF_W      = 210 - 2 * PDF_MARGIN

# Colour palette
_P  = (10,  132, 255)   # Blue  – primary
_G  = (48,  209,  88)   # Green – support
_O  = (255, 159,  10)   # Amber – extension / homework
_R  = (255,  69,  58)   # Red   – assess
_PU = (175, 100, 235)   # Purple – elaborate
_TL = (100, 210, 255)   # Teal  – explore
_DK = ( 29,  29,  31)   # Near-black text
_MD = (142, 142, 147)   # Mid grey
_LG = (242, 242, 247)   # Light grey fill
_WH = (255, 255, 255)

# 5E section colours: name-fragment → RGB
_5E_COLOURS = {
    "Engage":   (10,  132, 255),
    "Explore":  (100, 210, 255),
    "Explain":  (175, 100, 235),
    "Elaborate":(255, 159,  10),
    "Evaluate": ( 48, 209,  88),
}

def _5e_colour(section_name: str):
    for kw, col in _5E_COLOURS.items():
        if kw.lower() in section_name.lower():
            return col
    return _P


class LessonPlanPDF(FPDF):
    def __init__(self):
        super().__init__(orientation="P", unit="mm", format="A4")
        self.set_margins(PDF_MARGIN, PDF_MARGIN, PDF_MARGIN)
        self.set_auto_page_break(auto=True, margin=20)

    def header(self):
        pass

    def footer(self):
        self.set_y(-14)
        self.set_font("Helvetica", "", 7.5)
        self.set_text_color(*_MD)
        self.set_fill_color(*_LG)
        self.rect(0, self.get_y() - 1, 210, 16, "F")
        self.set_xy(PDF_MARGIN, self.get_y())
        self.cell(PDF_W / 2, 6, "Generated by EduCraft AI  |  Powered by Groq AI")
        self.cell(PDF_W / 2, 6, f"Page {self.page_no()}", align="R")


# ── Helper: coloured pill label ────────────────────────────────────────────────
def _pill(pdf, x, y, text, bg_rgb, text_rgb=_WH, font_size=7.5):
    pdf.set_font("Helvetica", "B", font_size)
    w = pdf.get_string_width(text) + 6
    pdf.set_fill_color(*bg_rgb)
    pdf.set_draw_color(*bg_rgb)
    pdf.rect(x, y, w, 5.5, "F")
    pdf.set_text_color(*text_rgb)
    pdf.set_xy(x + 1, y + 0.5)
    pdf.cell(w - 2, 4.5, text)
    return w


# ── Helper: section divider with coloured left bar ─────────────────────────────
def _section_header(pdf, name, dur, bloom, colour):
    row_h = 9
    pdf.set_fill_color(*_LG)
    pdf.rect(PDF_MARGIN, pdf.get_y(), PDF_W, row_h, "F")
    # left accent bar
    pdf.set_fill_color(*colour)
    pdf.rect(PDF_MARGIN, pdf.get_y(), 3, row_h, "F")

    # section name
    pdf.set_font("Helvetica", "B", 10.5)
    pdf.set_text_color(*_DK)
    pdf.set_xy(PDF_MARGIN + 5, pdf.get_y() + 1.5)
    pdf.cell(PDF_W * 0.52, 6, _safe(name))

    # duration pill
    dur_x = PDF_MARGIN + PDF_W * 0.55
    _pill(pdf, dur_x, pdf.get_y() + 1.5, f"{dur} min", colour)

    # bloom tag (right-aligned)
    pdf.set_font("Helvetica", "I", 8)
    pdf.set_text_color(*_MD)
    pdf.set_xy(PDF_MARGIN + PDF_W * 0.68, pdf.get_y() + 1.5)
    pdf.cell(PDF_W * 0.32 - 2, 6, _safe(bloom), align="R")

    pdf.ln(row_h + 2)


# ── Cover / header block ───────────────────────────────────────────────────────
def _pdf_cover(pdf, data):
    # Deep blue header band
    pdf.set_fill_color(*_P)
    pdf.rect(0, 0, 210, 38, "F")

    # Diagonal accent stripe (right side)
    pdf.set_fill_color(0, 90, 200)
    # simple rectangle stripe instead of polygon (fpdf2 basic)
    pdf.rect(160, 0, 50, 38, "F")

    pdf.set_xy(PDF_MARGIN, 6)
    pdf.set_font("Helvetica", "B", 18)
    pdf.set_text_color(*_WH)
    title = _safe(data.get("lesson_title", "Lesson Plan"))
    pdf.multi_cell(PDF_W - 20, 8, title)

    pdf.set_x(PDF_MARGIN)
    pdf.set_font("Helvetica", "", 9.5)
    pdf.set_text_color(200, 220, 255)
    meta = (f"{_safe(data.get('subject',''))}   |   {_safe(data.get('grade',''))}   |   "
            f"{data.get('duration_minutes','')} minutes")
    pdf.cell(0, 6, meta)
    pdf.ln(10)

    # ── Stats row ──────────────────────────────────────────────────────────────
    stats = [
        ("Sections",   str(len(data.get("sections", [])))),
        ("Objectives", str(len(data.get("learning_objectives", [])))),
        ("Materials",  str(len(data.get("materials", [])))),
        ("Framework",  "5E + Bloom"),
    ]
    box_w = PDF_W / len(stats)
    stat_y = pdf.get_y()
    for i, (label, val) in enumerate(stats):
        bx = PDF_MARGIN + i * box_w
        pdf.set_fill_color(*_LG)
        pdf.rect(bx, stat_y, box_w - 2, 14, "F")
        pdf.set_xy(bx + 2, stat_y + 1)
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(*_P)
        pdf.cell(box_w - 4, 7, val)
        pdf.set_xy(bx + 2, stat_y + 8)
        pdf.set_font("Helvetica", "", 7)
        pdf.set_text_color(*_MD)
        pdf.cell(box_w - 4, 5, label)
    pdf.ln(18)


# ── Objectives + Materials side by side ───────────────────────────────────────
def _pdf_obj_materials(pdf, data):
    col_w = PDF_W / 2 - 3
    right_x = PDF_MARGIN + col_w + 6

    # — Objectives column header
    _pill(pdf, PDF_MARGIN, pdf.get_y(), "LEARNING OBJECTIVES", _P, font_size=8)
    obj_start_y = pdf.get_y() + 7
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(*_DK)
    for i, obj in enumerate(data.get("learning_objectives", []), 1):
        pdf.set_xy(PDF_MARGIN, obj_start_y if i == 1 else pdf.get_y())
        # numbered circle
        pdf.set_fill_color(*_P)
        pdf.set_font("Helvetica", "B", 7.5)
        pdf.set_text_color(*_WH)
        pdf.rect(PDF_MARGIN, pdf.get_y() + 0.5, 5, 5, "F")
        pdf.set_xy(PDF_MARGIN + 0.5, pdf.get_y() + 0.5)
        pdf.cell(4, 4.5, str(i), align="C")
        # text
        pdf.set_font("Helvetica", "", 8.5)
        pdf.set_text_color(*_DK)
        pdf.set_xy(PDF_MARGIN + 6, pdf.get_y())
        pdf.multi_cell(col_w - 6, 5, _safe(obj))
        pdf.ln(1)
    obj_end_y = pdf.get_y()

    # — Materials column header
    pdf.set_y(pdf.get_y() - (obj_end_y - obj_start_y) - 7)  # reset to start
    _pill(pdf, right_x, pdf.get_y(), "MATERIALS & RESOURCES", _TL, font_size=8)
    mat_start_y = pdf.get_y() + 7
    pdf.set_font("Helvetica", "", 8.5)
    pdf.set_text_color(*_DK)
    for mat in data.get("materials", []):
        pdf.set_xy(right_x + 2, mat_start_y if mat == data["materials"][0] else pdf.get_y())
        pdf.set_fill_color(*_TL)
        pdf.rect(right_x + 2, pdf.get_y() + 1.5, 2, 2, "F")
        pdf.set_xy(right_x + 6, pdf.get_y())
        pdf.multi_cell(col_w - 6, 5, _safe(mat))
        pdf.ln(1)
    mat_end_y = pdf.get_y()

    pdf.set_y(max(obj_end_y, mat_end_y) + 5)


# ── Time-allocation bar ────────────────────────────────────────────────────────
def _pdf_timeline(pdf, sections):
    total = sum(s.get("duration_minutes", 0) for s in sections)
    if total == 0:
        return

    pdf.set_font("Helvetica", "B", 8)
    pdf.set_text_color(*_MD)
    pdf.cell(0, 5, "TIME ALLOCATION", ln=True)
    pdf.ln(1)

    bar_y = pdf.get_y()
    bar_h = 7
    x = PDF_MARGIN
    for section in sections:
        dur = section.get("duration_minutes", 0)
        col  = _5e_colour(section.get("name", ""))
        seg_w = (dur / total) * PDF_W
        pdf.set_fill_color(*col)
        pdf.rect(x, bar_y, seg_w, bar_h, "F")
        if seg_w > 8:
            pdf.set_font("Helvetica", "B", 6.5)
            pdf.set_text_color(*_WH)
            pdf.set_xy(x + 1, bar_y + 1.5)
            pdf.cell(seg_w - 2, 4, f"{dur}m", align="C")
        x += seg_w

    # Legend
    pdf.set_y(bar_y + bar_h + 3)
    x = PDF_MARGIN
    for section in sections:
        col  = _5e_colour(section.get("name", ""))
        name = section.get("name", "").split("(")[0].strip()
        short = name[:12]
        pdf.set_fill_color(*col)
        pdf.rect(x, pdf.get_y() + 1.5, 3.5, 3.5, "F")
        pdf.set_font("Helvetica", "", 6.5)
        pdf.set_text_color(*_DK)
        pdf.set_xy(x + 5, pdf.get_y())
        label_w = pdf.get_string_width(short) + 10
        pdf.cell(label_w, 6, short)
        x += label_w + 2
        if x > PDF_MARGIN + PDF_W - 30:
            pdf.ln(7)
            x = PDF_MARGIN

    pdf.ln(9)


# ── Individual section block ───────────────────────────────────────────────────
def _pdf_section(pdf, section):
    name   = section.get("name", "")
    dur    = section.get("duration_minutes", "")
    bloom  = section.get("bloom_level", "")
    colour = _5e_colour(name)

    _section_header(pdf, name, dur, bloom, colour)

    # Activities
    pdf.set_font("Helvetica", "B", 8)
    pdf.set_text_color(*colour)
    pdf.set_x(PDF_MARGIN + 3)
    pdf.cell(0, 5, "ACTIVITIES", ln=True)

    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(*_DK)
    for act in section.get("activities", []):
        pdf.set_x(PDF_MARGIN + 5)
        pdf.set_fill_color(*colour)
        pdf.rect(PDF_MARGIN + 5, pdf.get_y() + 2, 2, 2, "F")
        pdf.set_xy(PDF_MARGIN + 9, pdf.get_y())
        pdf.multi_cell(PDF_W - 9, 5, _safe(act))
    pdf.ln(2)

    # Teacher | Students columns
    col_w  = PDF_W / 2 - 4
    right_x = PDF_MARGIN + col_w + 8
    ty = pdf.get_y()

    # Teacher box
    pdf.set_fill_color(245, 248, 255)
    pdf.rect(PDF_MARGIN, ty, col_w, 6, "F")
    pdf.set_xy(PDF_MARGIN + 2, ty + 1)
    pdf.set_font("Helvetica", "B", 8)
    pdf.set_text_color(*_P)
    pdf.cell(col_w - 2, 4, "TEACHER ACTIONS")
    pdf.set_xy(PDF_MARGIN + 2, ty + 7)
    pdf.set_font("Helvetica", "", 8.5)
    pdf.set_text_color(*_DK)
    pdf.multi_cell(col_w - 2, 5, _safe(section.get("teacher_actions", "")))
    t_end = pdf.get_y()

    # Students box
    pdf.set_y(ty)
    pdf.set_fill_color(245, 255, 250)
    pdf.rect(right_x, ty, col_w, 6, "F")
    pdf.set_xy(right_x + 2, ty + 1)
    pdf.set_font("Helvetica", "B", 8)
    pdf.set_text_color(*_G)
    pdf.cell(col_w - 2, 4, "STUDENT ACTIONS")
    pdf.set_xy(right_x + 2, ty + 7)
    pdf.set_font("Helvetica", "", 8.5)
    pdf.set_text_color(*_DK)
    pdf.multi_cell(col_w - 2, 5, _safe(section.get("student_actions", "")))
    s_end = pdf.get_y()

    pdf.set_y(max(t_end, s_end) + 5)


# ── Differentiation ────────────────────────────────────────────────────────────
def _pdf_differentiation(pdf, diff):
    pdf.set_font("Helvetica", "B", 10)
    pdf.set_text_color(*_DK)
    # Left bar
    pdf.set_fill_color(*_O)
    pdf.rect(PDF_MARGIN, pdf.get_y(), 3, 8, "F")
    pdf.set_xy(PDF_MARGIN + 5, pdf.get_y() + 1)
    pdf.cell(0, 6, "DIFFERENTIATION STRATEGIES", ln=True)
    pdf.ln(2)

    col_w  = PDF_W / 2 - 4
    right_x = PDF_MARGIN + col_w + 8

    # Support header
    _pill(pdf, PDF_MARGIN, pdf.get_y(), "SUPPORT  (Struggling Learners)", _G, font_size=8)
    _pill(pdf, right_x, pdf.get_y(), "EXTENSION  (Advanced Learners)", _O, font_size=8)
    pdf.ln(7)

    sup_y = pdf.get_y()
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(*_DK)
    for s in diff.get("support", []):
        pdf.set_x(PDF_MARGIN + 3)
        pdf.set_fill_color(*_G)
        pdf.rect(PDF_MARGIN + 3, pdf.get_y() + 2, 2, 2, "F")
        pdf.set_xy(PDF_MARGIN + 7, pdf.get_y())
        pdf.multi_cell(col_w - 7, 5, _safe(s))
        pdf.ln(1)
    sup_end = pdf.get_y()

    pdf.set_y(sup_y)
    for e in diff.get("extension", []):
        pdf.set_x(right_x + 3)
        pdf.set_fill_color(*_O)
        pdf.rect(right_x + 3, pdf.get_y() + 2, 2, 2, "F")
        pdf.set_xy(right_x + 7, pdf.get_y())
        pdf.multi_cell(col_w - 7, 5, _safe(e))
        pdf.ln(1)
    ext_end = pdf.get_y()

    pdf.set_y(max(sup_end, ext_end) + 4)


# ── Homework ───────────────────────────────────────────────────────────────────
def _pdf_homework(pdf, homework):
    if not homework:
        return
    # Amber banner
    pdf.set_fill_color(255, 249, 235)
    hw_y = pdf.get_y()
    pdf.rect(PDF_MARGIN, hw_y, PDF_W, 5, "F")
    pdf.set_fill_color(*_O)
    pdf.rect(PDF_MARGIN, hw_y, 3, 5, "F")
    pdf.set_xy(PDF_MARGIN + 5, hw_y + 0.5)
    pdf.set_font("Helvetica", "B", 8.5)
    pdf.set_text_color(*_O)
    pdf.cell(0, 4, "OPTIONAL HOMEWORK")
    pdf.ln(6)
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(*_DK)
    pdf.set_x(PDF_MARGIN + 3)
    pdf.multi_cell(PDF_W - 3, 5, _safe(homework))
    pdf.ln(3)


# ── Public entry point ─────────────────────────────────────────────────────────
def create_lesson_plan_pdf(data: dict) -> bytes:
    pdf = LessonPlanPDF()
    pdf.add_page()

    _pdf_cover(pdf, data)
    _pdf_obj_materials(pdf, data)
    _pdf_timeline(pdf, data.get("sections", []))

    # Section heading
    pdf.set_fill_color(*_P)
    pdf.rect(PDF_MARGIN, pdf.get_y(), 3, 7, "F")
    pdf.set_xy(PDF_MARGIN + 5, pdf.get_y() + 1)
    pdf.set_font("Helvetica", "B", 11)
    pdf.set_text_color(*_DK)
    pdf.cell(0, 5, "LESSON SECTIONS  (5E Instructional Model)", ln=True)
    pdf.ln(3)

    for section in data.get("sections", []):
        _pdf_section(pdf, section)

    _pdf_differentiation(pdf, data.get("differentiation", {}))
    _pdf_homework(pdf, data.get("homework"))

    return bytes(pdf.output())


# ══════════════════════════════════════════════════════════════════════════════
#  QUESTION PAPER PDF  — student version + answer key
# ══════════════════════════════════════════════════════════════════════════════

# Colour palette for QP
_QP_BG   = (245, 245, 247)   # light grey page tint
_QP_HDR  = (10,  132, 255)   # blue header
_QP_MCQ  = (10,  132, 255)   # Section A
_QP_SHT  = (48,  209,  88)   # Section B
_QP_LNG  = (255, 159,  10)   # Section C
_QP_ANS  = (175, 100, 235)   # Answer key accent
_QP_DK   = (29,   29,  31)
_QP_MD   = (142, 142, 147)
_QP_WH   = (255, 255, 255)

_SECTION_COLOURS = {"mcq": _QP_MCQ, "short": _QP_SHT, "long": _QP_LNG}


class QuestionPaperPDF(FPDF):
    def __init__(self, is_answer_key=False):
        super().__init__(orientation="P", unit="mm", format="A4")
        self.set_margins(14, 14, 14)
        self.set_auto_page_break(auto=True, margin=18)
        self._is_answer_key = is_answer_key

    def header(self):
        pass

    def footer(self):
        self.set_y(-13)
        self.set_font("Helvetica", "", 7.5)
        self.set_text_color(*_QP_MD)
        label = "ANSWER KEY — FOR TEACHER USE ONLY" if self._is_answer_key else "Generated by EduCraft AI"
        self.cell(PDF_W / 2, 5, label)
        self.cell(PDF_W / 2, 5, f"Page {self.page_no()}", align="R")


def _qp_header(pdf, data, is_answer_key=False):
    """Big blue header block with paper title and meta."""
    col = _QP_ANS if is_answer_key else _QP_HDR

    pdf.set_fill_color(*col)
    pdf.rect(0, 0, 210, 36, "F")

    # Diagonal accent stripe
    pdf.set_fill_color(*[max(0, c - 40) for c in col])
    pdf.rect(155, 0, 55, 36, "F")

    pdf.set_xy(14, 5)
    pdf.set_font("Helvetica", "B", 16)
    pdf.set_text_color(*_QP_WH)
    label = "ANSWER KEY" if is_answer_key else data.get("board", ""),
    if is_answer_key:
        pdf.cell(0, 7, "ANSWER KEY & MARKING SCHEME", ln=True)
    else:
        pdf.cell(0, 7, _safe(f"{data.get('board','')} — {data.get('grade','')} {data.get('subject','')}"), ln=True)

    pdf.set_x(14)
    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 6, _safe(data.get("paper_title", data.get("subject", ""))), ln=True)

    pdf.set_x(14)
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(200, 220, 255)
    marks = data.get("total_marks", "")
    mins  = data.get("duration_minutes", "")
    diff  = data.get("difficulty", "")
    pdf.cell(0, 5, f"Total Marks: {marks}   |   Time: {mins} minutes   |   Difficulty: {diff}", ln=True)
    pdf.ln(8)


def _qp_instructions(pdf, instructions):
    """General instructions box."""
    if not instructions:
        return
    y = pdf.get_y()
    pdf.set_fill_color(*_QP_BG)
    # Draw box (estimate height)
    box_h = 6 + len(instructions) * 5.5
    pdf.rect(14, y, PDF_W, box_h, "F")
    pdf.set_fill_color(*_QP_HDR)
    pdf.rect(14, y, 3, box_h, "F")

    pdf.set_xy(20, y + 2)
    pdf.set_font("Helvetica", "B", 8.5)
    pdf.set_text_color(*_QP_DK)
    pdf.cell(0, 5, "GENERAL INSTRUCTIONS", ln=True)

    pdf.set_font("Helvetica", "", 8.5)
    for i, inst in enumerate(instructions, 1):
        pdf.set_x(20)
        pdf.cell(6, 5, f"{i}.")
        pdf.set_x(26)
        pdf.multi_cell(PDF_W - 14, 5, _safe(inst))
    pdf.ln(4)


def _qp_section_header(pdf, section, colour):
    """Coloured section header bar."""
    name  = section.get("section_name", "")
    mppq  = section.get("marks_per_question", 1)
    total = section.get("total_marks", 0)
    count = len(section.get("questions", []))
    stype = section.get("section_type", "")

    type_label = {"mcq": "Multiple Choice", "short": "Short Answer", "long": "Long Answer"}.get(stype, "")

    h = 9
    pdf.set_fill_color(*_QP_BG)
    pdf.rect(14, pdf.get_y(), PDF_W, h, "F")
    pdf.set_fill_color(*colour)
    pdf.rect(14, pdf.get_y(), 3, h, "F")

    pdf.set_xy(19, pdf.get_y() + 1.5)
    pdf.set_font("Helvetica", "B", 10)
    pdf.set_text_color(*_QP_DK)
    pdf.cell(PDF_W * 0.4, 6, _safe(f"{name} — {type_label}"))

    pdf.set_font("Helvetica", "", 8.5)
    pdf.set_text_color(*_QP_MD)
    pdf.cell(PDF_W * 0.35, 6, f"{count} questions  x  {mppq} mark{'s' if mppq>1 else ''} each")

    pdf.set_font("Helvetica", "B", 9)
    pdf.set_text_color(*colour)
    pdf.cell(PDF_W * 0.20, 6, f"{total} marks", align="R")
    pdf.ln(h + 3)


def _qp_mcq_student(pdf, question, colour):
    """Single MCQ — student version (question + options, no answer)."""
    num  = question.get("number", "")
    q    = _safe(question.get("question", ""))
    opts = question.get("options", [])

    # Question line
    pdf.set_font("Helvetica", "B", 9.5)
    pdf.set_text_color(*colour)
    pdf.set_x(14)
    pdf.cell(8, 5.5, f"Q{num}.")
    pdf.set_font("Helvetica", "", 9.5)
    pdf.set_text_color(*_QP_DK)
    pdf.set_x(22)
    pdf.multi_cell(PDF_W - 8, 5.5, q)

    # Options in 2 columns
    opt_w = PDF_W / 2
    for i, opt in enumerate(opts):
        col_x = 14 + (i % 2) * opt_w + 8
        if i % 2 == 0:
            oy = pdf.get_y()
        pdf.set_xy(col_x, oy if i % 2 == 1 else pdf.get_y())
        pdf.set_font("Helvetica", "", 9)
        pdf.set_text_color(*_QP_DK)
        pdf.cell(opt_w - 10, 5, _safe(opt))
        if i % 2 == 1:
            pdf.ln(5)
    if len(opts) % 2 == 1:
        pdf.ln(5)
    pdf.ln(2)


def _qp_mcq_answer(pdf, question, colour):
    """Single MCQ — answer key version."""
    num    = question.get("number", "")
    q      = _safe(question.get("question", ""))
    ans    = question.get("correct_answer", "")
    expl   = _safe(question.get("answer_explanation", ""))

    pdf.set_font("Helvetica", "B", 9.5)
    pdf.set_text_color(*colour)
    pdf.set_x(14)
    pdf.cell(8, 5.5, f"Q{num}.")
    pdf.set_font("Helvetica", "", 9.5)
    pdf.set_text_color(*_QP_DK)
    pdf.set_x(22)
    pdf.multi_cell(PDF_W - 8, 5.5, q)

    # Answer highlight
    pdf.set_x(22)
    pdf.set_fill_color(*[min(255, c + 190) for c in colour])
    pdf.set_font("Helvetica", "B", 9)
    pdf.set_text_color(*colour)
    pdf.cell(PDF_W - 8, 5, f"Answer: {ans}", fill=True)
    pdf.ln(5.5)

    if expl:
        pdf.set_x(22)
        pdf.set_font("Helvetica", "I", 8.5)
        pdf.set_text_color(*_QP_MD)
        pdf.multi_cell(PDF_W - 8, 4.5, f"Explanation: {expl}")
    pdf.ln(2)


def _qp_written_student(pdf, question, colour):
    """Short/long question — student version."""
    num   = question.get("number", "")
    q     = _safe(question.get("question", ""))
    bloom = question.get("bloom_level", "")
    marks = question.get("marks", "")

    pdf.set_font("Helvetica", "B", 9.5)
    pdf.set_text_color(*colour)
    pdf.set_x(14)
    pdf.cell(8, 5.5, f"Q{num}.")
    pdf.set_font("Helvetica", "", 9.5)
    pdf.set_text_color(*_QP_DK)
    pdf.set_x(22)
    pdf.multi_cell(PDF_W - 25, 5.5, q)

    # Bloom tag right-aligned
    if bloom:
        pdf.set_x(14)
        pdf.set_font("Helvetica", "I", 7.5)
        pdf.set_text_color(*_QP_MD)
        pdf.cell(PDF_W, 4, f"[{bloom}]", align="R")
    pdf.ln(3)


def _qp_written_answer(pdf, question, colour, marks_per_q):
    """Short/long question — answer key version."""
    num    = question.get("number", "")
    q      = _safe(question.get("question", ""))
    answer = _safe(question.get("model_answer", ""))
    scheme = question.get("marking_scheme", [])
    bloom  = question.get("bloom_level", "")

    pdf.set_font("Helvetica", "B", 9.5)
    pdf.set_text_color(*colour)
    pdf.set_x(14)
    pdf.cell(8, 5.5, f"Q{num}.")
    pdf.set_font("Helvetica", "", 9.5)
    pdf.set_text_color(*_QP_DK)
    pdf.set_x(22)
    pdf.multi_cell(PDF_W - 8, 5.5, q)

    # Model answer box
    pdf.set_x(22)
    pdf.set_fill_color(*_QP_BG)
    pdf.set_font("Helvetica", "B", 8)
    pdf.set_text_color(*colour)
    pdf.cell(PDF_W - 8, 5, "MODEL ANSWER:", fill=True)
    pdf.ln(5)
    pdf.set_x(22)
    pdf.set_font("Helvetica", "", 8.5)
    pdf.set_text_color(*_QP_DK)
    pdf.multi_cell(PDF_W - 8, 5, answer)

    # Marking scheme
    if scheme:
        pdf.ln(1)
        pdf.set_x(22)
        pdf.set_font("Helvetica", "B", 8)
        pdf.set_text_color(*_QP_ANS)
        pdf.cell(PDF_W - 8, 5, f"MARKING SCHEME  ({marks_per_q} marks):")
        pdf.ln(5)
        for point in scheme:
            pdf.set_x(26)
            pdf.set_fill_color(*_QP_ANS)
            pdf.rect(26, pdf.get_y() + 2, 2, 2, "F")
            pdf.set_xy(30, pdf.get_y())
            pdf.set_font("Helvetica", "", 8.5)
            pdf.set_text_color(*_QP_DK)
            pdf.multi_cell(PDF_W - 16, 4.5, _safe(point))

    if bloom:
        pdf.set_x(14)
        pdf.set_font("Helvetica", "I", 7.5)
        pdf.set_text_color(*_QP_MD)
        pdf.cell(PDF_W, 4, f"[{bloom}]", align="R")
    pdf.ln(4)


def _build_qp_pdf(data: dict, is_answer_key: bool) -> bytes:
    """Build either student or answer-key PDF from the same data dict."""
    pdf = QuestionPaperPDF(is_answer_key=is_answer_key)
    pdf.add_page()

    _qp_header(pdf, data, is_answer_key=is_answer_key)

    if not is_answer_key:
        _qp_instructions(pdf, data.get("general_instructions", []))

    for section in data.get("sections", []):
        stype  = section.get("section_type", "mcq")
        colour = _SECTION_COLOURS.get(stype, _QP_HDR)
        mppq   = section.get("marks_per_question", 1)

        _qp_section_header(pdf, section, colour)

        for question in section.get("questions", []):
            # Add marks_per_question to question dict for answer functions
            question.setdefault("marks", mppq)

            if is_answer_key:
                if stype == "mcq":
                    _qp_mcq_answer(pdf, question, colour)
                else:
                    _qp_written_answer(pdf, question, colour, mppq)
            else:
                if stype == "mcq":
                    _qp_mcq_student(pdf, question, colour)
                else:
                    _qp_written_student(pdf, question, colour)

        pdf.ln(4)

    return bytes(pdf.output())


def create_question_paper_pdfs(data: dict) -> tuple[bytes, bytes]:
    """Return (student_pdf_bytes, answer_key_pdf_bytes)."""
    student_pdf    = _build_qp_pdf(data, is_answer_key=False)
    answer_key_pdf = _build_qp_pdf(data, is_answer_key=True)
    return student_pdf, answer_key_pdf
